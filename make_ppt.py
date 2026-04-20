from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── PROFESSIONAL WHITE THEME COLORS ────────────────────
BG_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
PRIMARY    = RGBColor(0x1B, 0x3A, 0x5C)   # navy blue
ACCENT     = RGBColor(0x00, 0x7A, 0xCC)   # bright blue
ACCENT2    = RGBColor(0x00, 0x96, 0x88)   # teal
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)  # near-black
BODY_TEXT  = RGBColor(0x33, 0x33, 0x33)   # dark gray body
SUBTITLE   = RGBColor(0x66, 0x66, 0x66)   # medium gray
LIGHT_BG   = RGBColor(0xF2, 0xF6, 0xFA)  # light blue-gray
CARD_BG    = RGBColor(0xF8, 0xF9, 0xFB)  # very light gray
BORDER     = RGBColor(0xDE, 0xE2, 0xE8)  # border gray
RED_WARN   = RGBColor(0xDC, 0x35, 0x45)  # warning red
GREEN_OK   = RGBColor(0x28, 0xA7, 0x45)  # success green
ORANGE     = RGBColor(0xF0, 0x93, 0x22)  # orange
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TBL_HEADER = RGBColor(0x1B, 0x3A, 0x5C)  # table header
TBL_ALT1   = RGBColor(0xFF, 0xFF, 0xFF)  
TBL_ALT2   = RGBColor(0xF2, 0xF6, 0xFA)  

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_WHITE

def add_top_bar(slide):
    """Navy bar at the top of content slides"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

def add_shape_rect(slide, left, top, w, h, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, w, h, text, size=18, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

def slide_title(slide, title, subtitle=None):
    add_top_bar(slide)
    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             title, size=34, color=PRIMARY, bold=True)
    accent_line(slide, Inches(0.8), Inches(1.05), Inches(2))
    if subtitle:
        add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 subtitle, size=15, color=SUBTITLE)

def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    ts = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = "Calibri"
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = BODY_TEXT
                p.alignment = PP_ALIGN.LEFT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TBL_HEADER
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TBL_ALT1 if r % 2 == 1 else TBL_ALT2
    return ts


# ═══════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Navy block on left
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5.5), prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = PRIMARY
shape.line.fill.background()

# Accent stripe
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), 0, Inches(0.08), prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text(slide, Inches(0.8), Inches(1.8), Inches(4.5), Inches(1.5),
         "Smart Classroom\nOccupancy System", size=40, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(3.5), Inches(4.5), Inches(0.6),
         "with RFID Attendance", size=24, color=RGBColor(0x80, 0xBB, 0xE8))

add_text(slide, Inches(0.8), Inches(4.5), Inches(4.5), Inches(0.8),
         "An IoT-Based Real-Time Monitoring\n& Automated Attendance Solution",
         size=15, color=RGBColor(0xA0, 0xC0, 0xDE))

add_text(slide, Inches(0.8), Inches(5.8), Inches(4.5), Inches(0.4),
         "IoT Laboratory  |  April 2026", size=14, color=RGBColor(0x80, 0xA0, 0xC0))

# Right side - tools
add_text(slide, Inches(6.5), Inches(2.5), Inches(6), Inches(0.5),
         "Technology Stack", size=22, color=PRIMARY, bold=True)
accent_line(slide, Inches(6.5), Inches(3.05), Inches(1.5))

tools = ["ESP32 Microcontroller", "Wokwi Simulator", "ThingSpeak Cloud", "Node-RED Automation", "Arduino Framework"]
for i, tool in enumerate(tools):
    y = Inches(3.3) + Inches(i * 0.55)
    add_shape_rect(slide, Inches(6.5), y, Inches(5.5), Inches(0.42), LIGHT_BG, BORDER)
    add_text(slide, Inches(6.8), y + Inches(0.05), Inches(5), Inches(0.35),
             f"▸  {tool}", size=15, color=BODY_TEXT)


# ═══════════════════════════════════════════════════════
# SLIDE 2: TABLE OF CONTENTS
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_title(slide, "Table of Contents")

toc = [
    ("01", "Problem Statement", ACCENT),
    ("02", "Proposed Solution", ACCENT),
    ("03", "Methodology & Data Flow", ACCENT),
    ("04", "Hardware & Sensors", ACCENT),
    ("05", "Simulation & Results", ACCENT),
    ("06", "Alert System", ACCENT),
    ("07", "Limitations", ACCENT),
    ("08", "Future Scope", ACCENT),
    ("09", "Conclusion", ACCENT),
]

for i, (num, title, clr) in enumerate(toc):
    col = i // 5
    row = i % 5
    x = Inches(1.5) + Inches(col * 5.5)
    y = Inches(1.8) + Inches(row * 1.0)

    add_shape_rect(slide, x, y, Inches(5), Inches(0.75), LIGHT_BG, BORDER)
    add_text(slide, x + Inches(0.15), y + Inches(0.12), Inches(0.6), Inches(0.5),
             num, size=22, color=ACCENT, bold=True)

    # vertical line separator
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.8), y + Inches(0.1), Pt(2), Inches(0.55))
    sep.fill.solid()
    sep.fill.fore_color.rgb = BORDER
    sep.line.fill.background()

    add_text(slide, x + Inches(1.0), y + Inches(0.15), Inches(3.8), Inches(0.5),
             title, size=18, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════
# SLIDE 3: PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_title(slide, "1. Problem Statement", "Why traditional classrooms need a smart upgrade")

problems = [
    ("Manual Attendance", "Wastes 5–10 minutes per lecture.\nProxy attendance widespread with no physical verification.", ACCENT),
    ("No Occupancy Data", "Administrators have zero real-time\nvisibility into room usage.", ACCENT2),
    ("Poor Air Quality", "Overcrowded rooms → CO₂ > 2500 ppm →\ndrowsiness & reduced concentration.", ORANGE),
    ("Fire Safety Risk", "Safety norms require exact occupancy count.\nImpossible with paper registers.", RED_WARN),
    ("No Analytics", "No historical data on attendance trends,\npeak hours, or room utilization.", PRIMARY),
]

for i, (title, desc, clr) in enumerate(problems):
    col = i % 3
    row = i // 3
    x = Inches(0.6) + Inches(col * 4.1)
    y = Inches(1.7) + Inches(row * 2.7)

    card = add_shape_rect(slide, x, y, Inches(3.8), Inches(2.3), CARD_BG, BORDER)

    # Color accent bar on top of card
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.8), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = clr
    bar.line.fill.background()

    add_text(slide, x + Inches(0.2), y + Inches(0.25), Inches(3.4), Inches(0.45),
             title, size=18, color=PRIMARY, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.4), Inches(1.3),
             desc, size=14, color=SUBTITLE)


# ═══════════════════════════════════════════════════════
# SLIDE 4: PROPOSED SOLUTION
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_title(slide, "2. Proposed Solution", "Smart Classroom Occupancy System — Feature Overview")

features = [
    ("RFID Attendance", "Students tap college ID on\nRC522 reader → instant log", ACCENT),
    ("Occupancy Tracking", "Real-time entry/exit counting\nvia RFID + PIR motion sensor", ACCENT2),
    ("Air Quality Monitor", "MQ-135 sensor measures\nCO₂/VOC levels continuously", ORANGE),
    ("Environmental Sensing", "DHT22 tracks temperature &\nhumidity in real-time", GREEN_OK),
    ("Cloud Dashboard", "ThingSpeak displays live graphs\naccessible from any device", ACCENT),
    ("Smart Alerts", "Auto buzzer + LCD warnings on\novercrowding or bad air", RED_WARN),
]

for i, (title, desc, clr) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.6) + Inches(col * 4.1)
    y = Inches(1.7) + Inches(row * 2.7)

    card = add_shape_rect(slide, x, y, Inches(3.8), Inches(2.3), CARD_BG, BORDER)

    # colored circle indicator
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.3), Inches(0.3), Inches(0.3))
    circle.fill.solid()
    circle.fill.fore_color.rgb = clr
    circle.line.fill.background()

    add_text(slide, x + Inches(0.65), y + Inches(0.25), Inches(3), Inches(0.45),
             title, size=17, color=PRIMARY, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.85), Inches(3.4), Inches(1.2),
             desc, size=14, color=SUBTITLE)


# ═══════════════════════════════════════════════════════
# SLIDE 5: USPs
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_title(slide, "Unique Selling Points")

usps = [
    ("Replaces manual attendance registers entirely", GREEN_OK),
    ("Detects overcrowding before it becomes a safety hazard", GREEN_OK),
    ("Auto-logs who entered, when, and for how long", GREEN_OK),
    ("Real-world compatible with actual college RFID ID cards", GREEN_OK),
    ("Uses 100% open-source tools (ThingSpeak, Node-RED, Arduino)", GREEN_OK),
    ("Low cost — entire hardware system under ₹1,500", GREEN_OK),
]

for i, (usp, clr) in enumerate(usps):
    y = Inches(1.7) + Inches(i * 0.85)
    card = add_shape_rect(slide, Inches(1.5), y, Inches(10), Inches(0.65), LIGHT_BG, BORDER)

    check = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.7), y + Inches(0.12), Inches(0.4), Inches(0.4))
    check.fill.solid()
    check.fill.fore_color.rgb = clr
    check.line.fill.background()
    add_text(slide, Inches(1.75), y + Inches(0.1), Inches(0.35), Inches(0.4),
             "✓", size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(2.3), y + Inches(0.12), Inches(9), Inches(0.45),
             usp, size=18, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════
# SLIDE 6: METHODOLOGY (screenshot placeholder)
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_title(slide, "3. Methodology & Data Flow", "System Architecture")

# Left: big placeholder for circuit screenshot
card = add_shape_rect(slide, Inches(0.6), Inches(1.7), Inches(7), Inches(5.3), LIGHT_BG, BORDER)
add_text(slide, Inches(1.5), Inches(3.8), Inches(5.5), Inches(0.8),
         "[ Insert Wokwi Circuit Screenshot Here ]",
         size=18, color=SUBTITLE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(1.8), Inches(6.5), Inches(0.4),
         "Wokwi Simulation Circuit", size=16, color=PRIMARY, bold=True)

# Right: data flow steps
add_text(slide, Inches(8.2), Inches(1.6), Inches(4), Inches(0.4),
         "Data Flow", size=20, color=PRIMARY, bold=True)
accent_line(slide, Inches(8.2), Inches(2.05), Inches(1.2))

steps = [
    ("1", "SENSE", "Sensors read every 2s", ACCENT),
    ("2", "PROCESS", "ESP32 checks thresholds", ACCENT2),
    ("3", "ACT", "LCD + buzzer respond", ORANGE),
    ("4", "TRANSMIT", "HTTPS to ThingSpeak (16s)", ACCENT),
    ("5", "VISUALIZE", "Live graphs on cloud", GREEN_OK),
    ("6", "ALERT", "Node-RED triggers alerts", RED_WARN),
]

for i, (num, step, desc, clr) in enumerate(steps):
    y = Inches(2.3) + Inches(i * 0.75)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.2), y, Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = clr
    circle.line.fill.background()
    add_text(slide, Inches(8.23), y + Inches(0.04), Inches(0.35), Inches(0.35),
             num, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(8.8), y, Inches(4), Inches(0.3),
             step, size=15, color=PRIMARY, bold=True)
    add_text(slide, Inches(8.8), y + Inches(0.3), Inches(4), Inches(0.3),
             desc, size=12, color=SUBTITLE)


# ═══════════════════════════════════════════════════════
# SLIDE 7: PROTOCOLS
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_title(slide, "Communication Protocols")

proto_data = [
    ["Protocol", "Connection", "Purpose"],
    ["SPI", "ESP32 ↔ RC522 RFID", "High-speed card reading (10 MHz)"],
    ["I²C", "ESP32 ↔ LCD 16×2", "Display status messages (2 wires)"],
    ["OneWire", "ESP32 ↔ DHT22", "Temperature & humidity data"],
    ["ADC", "ESP32 ↔ MQ-135", "Analog air quality reading (0–4095)"],
    ["HTTPS", "ESP32 → ThingSpeak", "Encrypted cloud data transmission"],
    ["GPIO", "ESP32 ↔ PIR, Buzzer", "Digital input/output signals"],
]

add_table(slide, Inches(1.5), Inches(1.7), Inches(10), Inches(4.5),
          len(proto_data), 3, proto_data,
          col_widths=[Inches(2), Inches(3.5), Inches(4.5)])


# ═══════════════════════════════════════════════════════
# SLIDE 8: HARDWARE
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_title(slide, "4. Hardware & Sensors", "Component Specifications")

hw_data = [
    ["#", "Component", "Model", "Role", "Interface"],
    ["1", "Microcontroller", "ESP32 DevKit V1", "Central processing + WiFi", "—"],
    ["2", "RFID Reader", "MFRC522 (RC522)", "Student ID card scanning", "SPI"],
    ["3", "RFID Tags", "13.56 MHz Mifare", "Student ID cards", "—"],
    ["4", "Motion Sensor", "HC-SR501 PIR", "Presence detection", "GPIO"],
    ["5", "Air Quality", "MQ-135", "CO₂ / VOC detection", "ADC"],
    ["6", "Temp & Humidity", "DHT22 / AM2302", "Environmental monitoring", "OneWire"],
    ["7", "Display", "LCD 16×2 I2C", "Local status messages", "I²C"],
    ["8", "Buzzer", "Passive Buzzer", "Audio alerts", "GPIO"],
]

add_table(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(5),
          len(hw_data), 5, hw_data,
          col_widths=[Inches(0.5), Inches(2.2), Inches(2.8), Inches(3.5), Inches(2.5)])


# ═══════════════════════════════════════════════════════
# SLIDE 9: PIN MAP
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_title(slide, "ESP32 Pin Mapping")

pin_data = [
    ["Component", "ESP32 Pin", "Protocol"],
    ["RFID SDA", "GPIO 5", "SPI"],
    ["RFID SCK", "GPIO 18", "SPI"],
    ["RFID MOSI", "GPIO 23", "SPI"],
    ["RFID MISO", "GPIO 19", "SPI"],
    ["RFID RST", "GPIO 4", "GPIO"],
    ["LCD SDA", "GPIO 21", "I²C"],
    ["LCD SCL", "GPIO 22", "I²C"],
    ["DHT22 Data", "GPIO 15", "OneWire"],
    ["PIR Output", "GPIO 27", "Digital"],
    ["MQ-135 (Analog)", "GPIO 34", "ADC"],
    ["Buzzer", "GPIO 26", "Digital"],
]

add_table(slide, Inches(0.8), Inches(1.5), Inches(7.5), Inches(5.6),
          len(pin_data), 3, pin_data,
          col_widths=[Inches(2.8), Inches(2.2), Inches(2.5)])

# Why ESP32 box on right
card = add_shape_rect(slide, Inches(8.8), Inches(1.5), Inches(4), Inches(5.5), LIGHT_BG, BORDER)
add_text(slide, Inches(9.1), Inches(1.7), Inches(3.5), Inches(0.4),
         "Why ESP32?", size=20, color=PRIMARY, bold=True)
accent_line(slide, Inches(9.1), Inches(2.15), Inches(1.2))

reasons = [
    ("Dual-core", "240 MHz Xtensa processor"),
    ("WiFi Built-in", "802.11 b/g/n + BLE"),
    ("34 GPIO Pins", "Enough for all sensors"),
    ("ADC Support", "Reads analog MQ-135 directly"),
    ("Low Power", "Deep sleep modes available"),
    ("Affordable", "₹350–500 per unit"),
]
for i, (title, desc) in enumerate(reasons):
    y = Inches(2.4) + Inches(i * 0.72)
    add_text(slide, Inches(9.1), y, Inches(3.5), Inches(0.3),
             f"▸  {title}", size=14, color=ACCENT, bold=True)
    add_text(slide, Inches(9.4), y + Inches(0.3), Inches(3.2), Inches(0.3),
             desc, size=12, color=SUBTITLE)


# ═══════════════════════════════════════════════════════
# SLIDE 10: SIMULATION & RESULTS
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_title(slide, "5. Simulation & Results", "Wokwi ESP32 Simulator + ThingSpeak Cloud Dashboard")

test_data = [
    ["Test Scenario", "Action in Simulator", "Expected Result"],
    ["Student Entry", "Tap RFID tag on reader", "LCD: 'Welcome, Anshu!' + 1 beep + occupancy +1"],
    ["Student Exit", "Tap same tag again", "LCD: 'Goodbye, Anshu!' + 2 beeps + occupancy −1"],
    ["Unknown Card", "Tap unregistered tag", "LCD: 'Access Denied' + 3 beeps"],
    ["Temperature Change", "Click DHT22, adjust slider", "Updated values on LCD + ThingSpeak graph"],
    ["Air Quality Spike", "Turn potentiometer high", "LCD warning + 5 beeps if ADC > 2500"],
    ["Motion Detection", "Click PIR sensor", "Motion field updates on ThingSpeak cloud"],
]

add_table(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(4.5),
          len(test_data), 3, test_data,
          col_widths=[Inches(2.5), Inches(3.5), Inches(5.5)])

add_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         "→  ThingSpeak live dashboard results on next slide",
         size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 11: THINGSPEAK SCREENSHOTS
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_title(slide, "ThingSpeak Cloud Dashboard — Live Results")

fields = [
    ("Field 1: Temperature (°C)", 0, 0),
    ("Field 2: Humidity (%)", 1, 0),
    ("Field 3: Air Quality (ADC)", 0, 1),
    ("Field 4: Occupancy (count)", 1, 1),
    ("Field 5: Motion (0/1)", 0, 2),
]

for title, col, row in fields:
    x = Inches(0.6) + Inches(col * 6.2)
    y = Inches(1.6) + Inches(row * 1.85)
    w = Inches(5.8)
    h = Inches(1.65)

    card = add_shape_rect(slide, x, y, w, h, LIGHT_BG, BORDER)
    add_text(slide, x + Inches(0.15), y + Inches(0.05), w - Inches(0.3), Inches(0.3),
             title, size=13, color=PRIMARY, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.5), w - Inches(0.6), Inches(0.8),
             "[ Insert ThingSpeak Screenshot ]",
             size=15, color=SUBTITLE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.3),
         "Right-click placeholder → Format Shape → Fill → Picture → paste ThingSpeak screenshots",
         size=11, color=SUBTITLE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 12: SERIAL MONITOR
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_title(slide, "Serial Monitor Output")

serial_output = (
    "=== Smart Classroom Occupancy System ===\n"
    "[RFID] RC522 ready\n"
    "[WiFi] Connecting..... Connected!\n"
    "[WiFi] IP: 10.10.0.2\n"
    "[SYSTEM] Ready!\n"
    "\n"
    "[SENSORS] T=45.1C H=75.0% AQ=120 PIR=no Occ=0\n"
    "[RFID] UID: E9:76:42:A3\n"
    "[RFID] Anshu ENTERED | Occupancy: 1\n"
    "[CLOUD] ThingSpeak entry #1\n"
    "[SENSORS] T=45.1C H=75.0% AQ=3850 PIR=YES Occ=1\n"
    "[ALERT] OVERCROWDING / BAD AIR!\n"
    "[CLOUD] ThingSpeak entry #2\n"
    "[RFID] UID: A1:B2:C3:D4\n"
    "[RFID] Rahul ENTERED | Occupancy: 2\n"
    "[CLOUD] ThingSpeak entry #3"
)

card = add_shape_rect(slide, Inches(1), Inches(1.6), Inches(11), Inches(5.4), RGBColor(0x1E, 0x1E, 0x2E), BORDER)
add_text(slide, Inches(1.3), Inches(1.8), Inches(10.5), Inches(5.0),
         serial_output, size=15, color=RGBColor(0x50, 0xFA, 0x7B), font="Consolas")


# ═══════════════════════════════════════════════════════
# SLIDE 13: ALERT SYSTEM
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_title(slide, "6. Alert System", "Multi-level alerts for safety and attendance")

alert_data = [
    ["Alert Type", "Trigger Condition", "System Response"],
    ["Entry Beep", "Valid RFID card tapped (entry)", "1 beep + LCD: 'Welcome, [Name]!'"],
    ["Exit Beep", "Same card tapped again (exit)", "2 beeps + LCD: 'Goodbye, [Name]!'"],
    ["Unknown Card", "Unregistered RFID tag detected", "3 beeps + LCD: 'Access Denied'"],
    ["Overcrowding", "Occupancy count > 30 students", "5 rapid beeps + LCD: '!! WARNING !!'"],
    ["Bad Air Quality", "MQ-135 analog > 2500 (ADC)", "5 rapid beeps + LCD: 'BAD AIR QUALITY'"],
]

add_table(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(3.5),
          len(alert_data), 3, alert_data,
          col_widths=[Inches(2.5), Inches(4), Inches(5)])

# Node-RED box
card = add_shape_rect(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.6), LIGHT_BG, BORDER)
add_text(slide, Inches(1.1), Inches(5.6), Inches(5), Inches(0.35),
         "Node-RED Extensions (Future)", size=16, color=PRIMARY, bold=True)
extensions = [
    "📧  Email notification when occupancy > threshold",
    "📱  SMS alert to faculty for dangerous CO₂",
    "🌀  Auto-trigger fans/AC via relay module",
    "📋  Log attendance to Google Sheets automatically",
]
for i, ext in enumerate(extensions):
    col = i % 2
    row = i // 2
    x = Inches(1.1) + Inches(col * 5.8)
    y = Inches(6.05) + Inches(row * 0.35)
    add_text(slide, x, y, Inches(5.5), Inches(0.3),
             ext, size=13, color=BODY_TEXT)


# ═══════════════════════════════════════════════════════
# SLIDE 14: LIMITATIONS
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_title(slide, "7. Limitations")

limitations = [
    ("ThingSpeak Rate Limit", "Free tier: updates only every 15 seconds — not truly real-time"),
    ("Single Entry Point", "Current design assumes one door; multi-door needs more readers"),
    ("No Persistent Database", "Attendance log in serial output only — no database storage"),
    ("MQ-135 Calibration", "Simulated with potentiometer; real sensor needs warm-up period"),
    ("WiFi Dependency", "Cloud features fail without internet (local LCD/buzzer still work)"),
    ("Hardcoded API Key", "Needs secure key management for production deployment"),
    ("No Battery Backup", "ESP32 needs constant 5V USB — no operation during power cuts"),
    ("Simulation vs Real HW", "Wokwi RFID has limited features vs actual RC522 hardware"),
]

for i, (title, desc) in enumerate(limitations):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + Inches(col * 6.2)
    y = Inches(1.6) + Inches(row * 1.35)

    card = add_shape_rect(slide, x, y, Inches(5.9), Inches(1.15), CARD_BG, BORDER)

    # Red left bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Pt(4), Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()

    add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.5), Inches(0.3),
             f"⚠  {title}", size=15, color=ORANGE, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.5), Inches(5.5), Inches(0.55),
             desc, size=13, color=SUBTITLE)


# ═══════════════════════════════════════════════════════
# SLIDE 15: FUTURE SCOPE
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_title(slide, "8. Future Scope", "Potential Enhancements & Scalability")

future = [
    ("Face Recognition", "ESP32-CAM for visual verification — zero proxy", ACCENT),
    ("Mobile App", "Flutter/React Native for student attendance check", ACCENT2),
    ("Database Backend", "Firebase/MongoDB for persistent attendance records", GREEN_OK),
    ("Multi-Door Support", "Multiple ESP32+RFID units per classroom", ACCENT),
    ("HVAC Integration", "Auto-control fans/AC based on temp + CO₂", ORANGE),
    ("ML Prediction", "Predict peak occupancy from historical data", ACCENT2),
    ("LoRa/Mesh Network", "Campus-wide coverage without WiFi", PRIMARY),
    ("Solar + Battery", "Uninterrupted operation during power cuts", GREEN_OK),
]

for i, (title, desc, clr) in enumerate(future):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + Inches(col * 6.2)
    y = Inches(1.6) + Inches(row * 1.35)

    card = add_shape_rect(slide, x, y, Inches(5.9), Inches(1.15), CARD_BG, BORDER)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Pt(4), Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = clr
    bar.line.fill.background()

    add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.5), Inches(0.3),
             f"→  {title}", size=15, color=PRIMARY, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.5), Inches(5.5), Inches(0.55),
             desc, size=13, color=SUBTITLE)


# ═══════════════════════════════════════════════════════
# SLIDE 16: CONCLUSION
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_title(slide, "9. Conclusion")

conclusions = [
    "Successfully designed and simulated a Smart Classroom Occupancy System using ESP32",
    "RFID-based attendance eliminates manual registers and prevents proxy attendance",
    "Real-time environmental monitoring (temp, humidity, CO₂) ensures student comfort",
    "Cloud dashboard on ThingSpeak provides live visualization from any device",
    "Alert system warns against overcrowding and poor air quality — locally and on cloud",
    "Entire system built using 100% open-source tools as per assignment requirements",
]

for i, c in enumerate(conclusions):
    y = Inches(1.6) + Inches(i * 0.8)
    card = add_shape_rect(slide, Inches(1.2), y, Inches(10.5), Inches(0.62), LIGHT_BG, BORDER)

    check = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.4), y + Inches(0.12), Inches(0.38), Inches(0.38))
    check.fill.solid()
    check.fill.fore_color.rgb = GREEN_OK
    check.line.fill.background()
    add_text(slide, Inches(1.43), y + Inches(0.1), Inches(0.35), Inches(0.38),
             "✓", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(2.0), y + Inches(0.12), Inches(9.5), Inches(0.45),
             c, size=16, color=DARK_TEXT)

# Key takeaway
card = add_shape_rect(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.7), ACCENT, None)
add_text(slide, Inches(1.8), Inches(6.55), Inches(9.5), Inches(0.6),
         "Low-cost IoT hardware + open-source cloud = smarter, safer classrooms",
         size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 17: OPEN SOURCE TOOLS
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_title(slide, "Open-Source Tools Used", "As per assignment requirements")

tools_data = [
    ["Tool", "Category", "Role in Project"],
    ["Arduino IDE", "Development Platform", "C++ firmware development for ESP32"],
    ["Wokwi", "Simulator", "Virtual ESP32 + sensor circuit simulation"],
    ["ThingSpeak", "Cloud Platform (MathWorks)", "Real-time dashboard & data storage"],
    ["Node-RED", "Automation Engine", "Alert rules & workflow automation"],
    ["MFRC522 Library", "Open-source Arduino Library", "RFID reader communication driver"],
    ["PlatformIO", "Build System", "Compilation & dependency management"],
]

add_table(slide, Inches(1.5), Inches(1.7), Inches(10), Inches(4.2),
          len(tools_data), 3, tools_data,
          col_widths=[Inches(3), Inches(3.5), Inches(3.5)])


# ═══════════════════════════════════════════════════════
# SLIDE 18: REFERENCES
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_title(slide, "References")

refs = [
    "[1]  Wokwi ESP32 Simulator — https://wokwi.com",
    "[2]  ThingSpeak IoT Platform — https://thingspeak.com",
    "[3]  Node-RED Flow Editor — https://nodered.org",
    "[4]  MFRC522 Library — https://github.com/miguelbalboa/rfid",
    "[5]  ESP32 Arduino Core — https://github.com/espressif/arduino-esp32",
    "[6]  Tech Monitor: Open-Source IoT Tools",
    "       https://techmonitor.ai/technology/software/open-source-tools-for-iot-applications",
    "[7]  G2: IoT Platforms — https://g2.com/categories/iot-platforms/free",
    "[8]  DHT22 Datasheet — Aosong Electronics",
    "[9]  MQ-135 Gas Sensor Datasheet — Winsen Electronics",
]

txBox = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(11), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True
for ref in refs:
    p = tf.add_paragraph()
    p.text = ref
    p.font.size = Pt(15)
    p.font.color.rgb = BODY_TEXT
    p.font.name = "Calibri"
    p.space_after = Pt(6)


# ═══════════════════════════════════════════════════════
# SLIDE 19: THANK YOU
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Navy block on the right
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), 0, Inches(5.833), prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = PRIMARY
shape.line.fill.background()

# Accent stripe
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.42), 0, Inches(0.08), prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text(slide, Inches(1), Inches(2.2), Inches(6), Inches(1),
         "Thank You!", size=52, color=PRIMARY, bold=True)

accent_line(slide, Inches(1), Inches(3.4), Inches(2))

add_text(slide, Inches(1), Inches(3.7), Inches(6), Inches(0.7),
         "Smart Classroom Occupancy System\nwith RFID Attendance",
         size=18, color=SUBTITLE)

add_text(slide, Inches(8), Inches(2.8), Inches(5), Inches(1),
         "Questions?", size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(8), Inches(4.5), Inches(5), Inches(1.5),
         "ESP32 · Wokwi · ThingSpeak\nNode-RED · Arduino",
         size=16, color=RGBColor(0x80, 0xBB, 0xE8), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════
output_path = r"c:\Users\anshu\Desktop\ADA2\iot\Smart_Classroom_PPT.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
